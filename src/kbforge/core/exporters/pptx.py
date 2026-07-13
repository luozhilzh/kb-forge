"""PPTX exporter — built-in minimal template (no external .pptx dependency).

Renders a :class:`CaseBundle` list to a PowerPoint file using ``python-pptx``.
The layout is generated in code (title slide + one slide per bundle), so there
is no template file to ship and CI can validate the output headlessly.

Per the chosen fork: built-in minimal template (fork A), deterministic, no
external model/vector dependency.
"""

from __future__ import annotations

from pathlib import Path

from .base import BaseExporter
from .extract import CaseBundle

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:  # pragma: no cover — python-pptx is an optional runtime dep
    Presentation = None


class PptxExporter(BaseExporter):
    name = "pptx"
    suffix = ".pptx"

    def export(self, pages: list[CaseBundle], out_path: Path) -> Path:
        if Presentation is None:
            raise RuntimeError(
                "python-pptx is required for the pptx exporter "
                "(pip install python-pptx)"
            )
        out_path = out_path.with_suffix(".pptx")
        prs = Presentation()

        # --- title slide ---
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = "kb-forge 案例导出"
        sub = title_slide.placeholders[1]
        sub.text = f"共 {len(pages)} 个条目 · 由知识库自动生成"

        # --- one slide per bundle ---
        for b in pages:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = b.title

            body = slide.placeholders[1].text_frame
            body.clear()
            p0 = body.paragraphs[0]
            p0.text = f"[{b.type}] {b.summary}" if b.summary else f"[{b.type}]"
            if b.content:
                p = body.add_paragraph()
                p.text = b.content
                p.font.size = Pt(14)
            if b.source_anchor:
                p = body.add_paragraph()
                p.text = "来源：" + "；".join(b.source_anchor)
                p.font.size = Pt(10)
            if b.related:
                p = body.add_paragraph()
                p.text = "相关：" + "，".join(f"[[{r}]]" for r in b.related)
                p.font.size = Pt(10)

        out_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(out_path)
        return out_path
