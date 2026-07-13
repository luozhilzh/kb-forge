"""Golden tests for the exporter family (extract -> pptx / html / md).

Verifies that a built wiki can be turned into structured CaseBundles and then
rendered to each format with assertions that survive headless CI.
"""

from __future__ import annotations

from pathlib import Path

from kbforge.core.exporters import (
    MarkdownExporter,
    PptxExporter,
    HtmlExporter,
    extract_bundles,
    CaseBundle,
)


def test_extract_bundles_from_built_wiki(built_wiki) -> None:
    bundles = extract_bundles(built_wiki)
    assert len(bundles) == 3
    by_id = {b.id: b for b in bundles}
    assert by_id["t101-rag-eval"].type == "case"
    assert by_id["t102-chunk-strategy"].type == "concept"
    assert by_id["t103-gpu-cost"].type == "pitfall"
    # related links are carried from [[wiki-links]] in the body
    assert "t102-chunk-strategy" in by_id["t101-rag-eval"].related
    assert "t103-gpu-cost" in by_id["t101-rag-eval"].related
    # source anchor preserved from the page frontmatter
    assert by_id["t101-rag-eval"].source_anchor
    # built-in minimal summary is the first paragraph
    assert by_id["t101-rag-eval"].summary
    # type filter works
    cases_only = extract_bundles(built_wiki, types=frozenset({"case"}))
    assert [b.id for b in cases_only] == ["t101-rag-eval"]


def test_pptx_exporter_roundtrip(built_wiki, tmp_path) -> None:
    bundles = extract_bundles(built_wiki)
    out = PptxExporter().export(bundles, tmp_path / "cases.pptx")

    assert out.exists() and out.suffix == ".pptx"
    from pptx import Presentation

    prs = Presentation(str(out))
    # title slide + one per bundle
    assert len(prs.slides) == len(bundles) + 1

    texts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
    joined = "\n".join(texts)
    for b in bundles:
        assert b.title in joined, f"title {b.title!r} missing from pptx"


def test_html_exporter_roundtrip(built_wiki, tmp_path) -> None:
    bundles = extract_bundles(built_wiki)
    out = HtmlExporter().export(bundles, tmp_path / "cases.html")

    assert out.exists() and out.suffix == ".html"
    html_text = out.read_text(encoding="utf-8")
    assert "<!DOCTYPE html>" in html_text
    for b in bundles:
        assert b.title in html_text
        assert f"data-id='{b.id}'" in html_text
        # source anchor rendered
        if b.source_anchor:
            assert b.source_anchor[0] in html_text


def test_markdown_exporter_still_works(built_wiki, tmp_path) -> None:
    bundles = extract_bundles(built_wiki)
    out = MarkdownExporter().export(bundles, tmp_path / "cases.md")
    assert out.exists() and out.suffix == ".md"
    content = out.read_text(encoding="utf-8")
    for b in bundles:
        assert b.title in content
